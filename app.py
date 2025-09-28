# app.py
# ✈️ VNA Tech – Tra cứu tài liệu từ Google Drive (PDF/PPTX)
# Tác vụ: Đồng bộ Drive -> chunking -> embedding -> cache (.pkl) -> truy hồi Top-K -> trả lời.

import os, json, hashlib, pickle, tempfile
from collections.abc import Mapping
from typing import Dict, Any, List, Tuple

import numpy as np
import streamlit as st
import tiktoken
from pypdf import PdfReader
from pptx import Presentation
from openai import OpenAI

from pydrive2.auth import GoogleAuth
from pydrive2.drive import GoogleDrive

# ================== Cấu hình chung ==================
APP_TITLE = "✈️ VNA Tech: Hỗ trợ tra cứu thông tin tài liệu từ Google Drive"
EMBED_MODEL = "text-embedding-3-small"   # 1536 chiều
EMBED_DIM = 1536
CHUNK_TOKENS = 400
CHUNK_OVERLAP = 80
TOP_K = 10
META_PATH = "embeddings_meta.pkl"        # cache chính (metadata + embeddings + texts)
# ====================================================


# ----------------- Helpers: Secrets -----------------
def require_secret(key: str) -> str:
    val = st.secrets.get(key) or os.getenv(key)
    if not val:
        st.error(f"Thiếu {key} trong Secrets."); st.stop()
    return val

def load_gsa_secrets() -> Dict[str, Any]:
    """
    Đọc GOOGLE_SERVICE_ACCOUNT_JSON dưới dạng:
    - Mapping (TOML object) -> trả về dict
    - String JSON (triple quotes '''...''' hoặc """...""") -> json.loads
    """
    raw = st.secrets.get("GOOGLE_SERVICE_ACCOUNT_JSON", None)
    if raw is None:
        st.error("Thiếu GOOGLE_SERVICE_ACCOUNT_JSON trong Secrets."); st.stop()

    if isinstance(raw, Mapping):
        return dict(raw)  # TOML object

    if isinstance(raw, str):
        try:
            return json.loads(raw)  # JSON string
        except Exception:
            st.error("GOOGLE_SERVICE_ACCOUNT_JSON không phải JSON hợp lệ. Kiểm tra triple quotes và giữ nguyên ký tự \\n trong private_key.")
            st.stop()

    st.error(f"GOOGLE_SERVICE_ACCOUNT_JSON có kiểu không hỗ trợ: {type(raw).__name__}")
    st.stop()


# ----------------- Google Drive -----------------
def get_drive(creds_dict: Dict[str, Any]) -> GoogleDrive:
    """Khởi tạo GoogleDrive dùng service account qua settings (truyền client_json là dict)."""
    gauth = GoogleAuth(settings={
        "client_config_backend": "service",
        "service_config": {"client_json": creds_dict},
        "save_credentials": False,
    })
    gauth.ServiceAuth()  # không truyền tham số -> dùng settings ở trên
    return GoogleDrive(gauth)

def list_drive_files(drive: GoogleDrive, folder_id: str) -> List[Dict[str, Any]]:
    # Hỗ trợ PDF và PPTX; bật include All Drives cho Shared Drive
    q = (
        f"'{folder_id}' in parents and trashed=false and ("
        "mimeType='application/pdf' or "
        "mimeType='application/vnd.openxmlformats-officedocument.presentationml.presentation'"
        ")"
    )
    files = drive.ListFile({
        "q": q,
        "supportsAllDrives": True,
        "includeItemsFromAllDrives": True
    }).GetList()
    return [
        {
            "id": f["id"],
            "title": f["title"],
            "mimeType": f["mimeType"],
            "modifiedDate": f.get("modifiedDate"),
            "md5Checksum": f.get("md5Checksum"),
        }
        for f in files
    ]

def download_drive_file(drive: GoogleDrive, file_id: str, local_path: str) -> str:
    f = drive.CreateFile({"id": file_id})
    f.GetContentFile(local_path)
    return local_path


# ----------------- Xử lý tài liệu & embedding -----------------
def md5_of_file(path: str) -> str:
    h = hashlib.md5()
    with open(path, "rb") as fp:
        for chunk in iter(lambda: fp.read(8192), b""):
            h.update(chunk)
    return h.hexdigest()

def read_pdf(path: str) -> List[Tuple[int, str]]:
    out = []
    reader = PdfReader(path)
    for i, p in enumerate(reader.pages, start=1):
        txt = p.extract_text() or ""
        out.append((i, txt))
    return out

def read_pptx(path: str) -> List[Tuple[int, str]]:
    prs = Presentation(path)
    out = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = [sh.text for sh in slide.shapes if hasattr(sh, "text")]
        out.append((i, " ".join(parts)))
    return out

def token_chunks(text: str, max_tokens=CHUNK_TOKENS, overlap=CHUNK_OVERLAP, enc_name="cl100k_base") -> List[str]:
    enc = tiktoken.get_encoding(enc_name)
    toks = enc.encode(text)
    chunks, i, n = [], 0, len(toks)
    while i < n:
        j = min(i + max_tokens, n)
        chunks.append(enc.decode(toks[i:j]))
        if j == n: break
        i = j - overlap
    return chunks

def get_openai_client() -> OpenAI:
    return OpenAI(api_key=require_secret("OPENAI_API_KEY"))

def embed_texts(client: OpenAI, texts: List[str]) -> np.ndarray:
    if not texts: return np.zeros((0, EMBED_DIM), dtype=np.float32)
    r = client.embeddings.create(model=EMBED_MODEL, input=texts)
    vecs = [d.embedding for d in r.data]
    return np.array(vecs, dtype=np.float32)

def l2_normalize(X: np.ndarray) -> np.ndarray:
    n = np.linalg.norm(X, axis=1, keepdims=True) + 1e-10
    return X / n


# ----------------- Cache embeddings_meta.pkl -----------------
# Cấu trúc:
# {
#   "dim": 1536,
#   "files": { <file_id>: {"title":..., "md5":..., "ranges":[(start,end)]} },
#   "embeddings": np.ndarray (N,dim)  -- đã chuẩn hoá,
#   "texts": List[str]  -- song hành với embeddings
# }
def load_meta(path=META_PATH):
    if os.path.exists(path):
        with open(path, "rb") as f: return pickle.load(f)
    return {"dim": EMBED_DIM, "files": {}, "embeddings": np.zeros((0, EMBED_DIM), dtype=np.float32), "texts": []}

def save_meta(meta, path=META_PATH):
    with open(path, "wb") as f: pickle.dump(meta, f)

def sync_from_drive(drive: GoogleDrive, folder_id: str):
    """Đồng bộ: chỉ xử lý file mới/đổi, sau đó cập nhật embeddings_meta.pkl"""
    client = get_openai_client()
    meta = load_meta()
    files = list_drive_files(drive, folder_id)
    added = 0

    for f in files:
        fid, title = f["id"], f["title"]

        with tempfile.TemporaryDirectory() as td:
            local = os.path.join(td, title)
            download_drive_file(drive, fid, local)
            md5_now = md5_of_file(local)

            # Bỏ qua nếu chưa thay đổi
            if fid in meta["files"] and meta["files"][fid]["md5"] == md5_now:
                continue

            # Đọc & chunk
            pages = read_pdf(local) if title.lower().endswith(".pdf") else read_pptx(local)
            chunks = []
            for _, txt in pages:
                chunks.extend(token_chunks(txt))

            # Embedding & chuẩn hoá
            vecs = embed_texts(client, chunks)
            vecs = l2_normalize(vecs)

            # Append vào cache
            start = meta["embeddings"].shape[0]
            if vecs.shape[0] > 0:
                meta["embeddings"] = np.vstack([meta["embeddings"], vecs])
                meta["texts"].extend(chunks)
            end = meta["embeddings"].shape[0]

            meta["files"][fid] = {"title": title, "md5": md5_now, "ranges": [(start, end)]}
            added += (end - start)

    if added > 0: save_meta(meta)
    return meta

def retrieve_top_k(query: str, meta, k=TOP_K) -> List[Tuple[float, str, str]]:
    client = get_openai_client()
    q = embed_texts(client, [query])
    q = l2_normalize(q)[0]  # (dim,)
    M = meta["embeddings"]
    if M.shape[0] == 0: return []

    sims = M @ q  # cosine (đã chuẩn hoá)
    idx = np.argsort(-sims)[:k]

    # ánh xạ index -> file title
    file_by_index = {}
    for fid, finfo in meta["files"].items():
        title = finfo["title"]
        for (s, e) in finfo["ranges"]:
            for i in range(s, e):
                file_by_index[i] = title

    results = []
    for i in idx:
        results.append((float(sims[i]), file_by_index.get(int(i), "unknown"), meta["texts"][int(i)]))
    return results

def answer_from_chunks(query: str, chunks: List[str]) -> str:
    client = get_openai_client()
    context = "\n\n".join(chunks)
    r = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {"role": "system", "content": "Bạn là trợ lý kỹ thuật. Trả lời ngắn gọn, chỉ dựa trên CONTEXT."},
            {"role": "user", "content": f"CONTEXT:\n{context}\n\nCÂU HỎI: {query}\n\nNếu thiếu dữ liệu, hãy nêu rõ."}
        ],
        temperature=0.2
    )
    return r.choices[0].message.content


# ----------------- UI -----------------
st.set_page_config(page_title="VNA Tech – Tra cứu tài liệu", page_icon="✈️", layout="wide")
st.title(APP_TITLE)

# Secrets tối thiểu
DRIVE_FOLDER_ID = require_secret("DRIVE_FOLDER_ID")
creds_dict = load_gsa_secrets()  # hỗ trợ cả string JSON lẫn object TOML

with st.sidebar:
    st.subheader("Thiết lập")
    st.write("• Nguồn: Google Drive (thư mục đã chia sẻ cho Service Account).")
    if st.button("🔄 Đồng bộ Drive → Cache (.pkl)"):
        with st.spinner("Đang đồng bộ dữ liệu từ Google Drive..."):
            drive = get_drive(creds_dict)
            meta_after = sync_from_drive(drive, DRIVE_FOLDER_ID)
            st.success(f"Đồng bộ xong. Tổng vector: {meta_after['embeddings'].shape[0]}")

st.markdown("—")
query = st.text_input("Đặt câu hỏi:")
if st.button("Tìm câu trả lời") and query:
    drive = get_drive(creds_dict)
    meta = sync_from_drive(drive, DRIVE_FOLDER_ID)  # chỉ xử lý file mới/đổi
    if meta["embeddings"].shape[0] == 0:
        st.warning("Chưa có dữ liệu trong cache. Hãy bấm 'Đồng bộ Drive → Cache (.pkl)'."); st.stop()

    hits = retrieve_top_k(query, meta, k=TOP_K)
    if not hits:
        st.info("Không tìm thấy đoạn liên quan."); st.stop()

    st.subheader("🔟 Các đoạn liên quan")
    chunks = []
    for score, title, chunk in hits:
        st.markdown(f"**[{title}]** — score={score:.3f}")
        st.write((chunk[:800] + "...") if len(chunk) > 800 else chunk)
        st.markdown("---")
        chunks.append(chunk)

    st.subheader("✨ Trả lời")
    st.write(answer_from_chunks(query, chunks))
