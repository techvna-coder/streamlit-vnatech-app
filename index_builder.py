# index_builder.py
import os, pickle, hashlib
import numpy as np
import faiss
import tiktoken
from pypdf import PdfReader
from pptx import Presentation
from openai import OpenAI

EMBED_MODEL = "text-embedding-3-small"
DIM = 1536  # chiều embedding của model trên

def get_client(key: str) -> OpenAI:
    return OpenAI(api_key=key)

def md5_bytes(b: bytes) -> str:
    import hashlib; return hashlib.md5(b).hexdigest()

def read_pdf(path: str) -> list[tuple[int, str]]:
    out = []
    reader = PdfReader(path)
    for i, p in enumerate(reader.pages, start=1):
        txt = p.extract_text() or ""
        out.append((i, txt))
    return out

def read_pptx(path: str) -> list[tuple[int, str]]:
    prs = Presentation(path)
    out = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = [sh.text for sh in slide.shapes if hasattr(sh, "text")]
        out.append((i, " ".join(parts)))
    return out

def chunk_text(text: str, max_tokens=400, overlap=80, encoding="cl100k_base"):
    enc = tiktoken.get_encoding(encoding)
    toks = enc.encode(text)
    i, n = 0, len(toks)
    while i < n:
        j = min(i + max_tokens, n)
        yield enc.decode(toks[i:j])
        if j == n: break
        i = j - overlap

def embed_texts(client: OpenAI, texts: list[str]) -> np.ndarray:
    if not texts: return np.zeros((0, DIM), dtype=np.float32)
    r = client.embeddings.create(model=EMBED_MODEL, input=texts)
    vecs = [d.embedding for d in r.data]
    return np.array(vecs, dtype=np.float32)

def l2_normalize(X: np.ndarray) -> np.ndarray:
    n = np.linalg.norm(X, axis=1, keepdims=True) + 1e-10
    return X / n

def load_meta(path="embeddings_meta.pkl"):
    if os.path.exists(path):
        with open(path, "rb") as f: return pickle.load(f)
    # cấu trúc meta:
    # {
    #   "files": { file_id: {"title":..., "md5":..., "modified":..., "chunks":[(start,end)]} },
    #   "dim": DIM, "index_path":"faiss_index.bin"
    # }
    return {"files": {}, "dim": DIM, "index_path": "faiss_index.bin"}

def save_meta(meta, path="embeddings_meta.pkl"):
    with open(path, "wb") as f: pickle.dump(meta, f)

def load_faiss(index_path: str):
    if os.path.exists(index_path):
        return faiss.read_index(index_path)
    return faiss.IndexFlatIP(DIM)  # dùng cosine (sau khi chuẩn hoá)

def save_faiss(index, index_path: str):
    faiss.write_index(index, index_path)

def add_document_to_index(client: OpenAI, index, meta, file_id, title, local_path, kind: str):
    # Trả về số vector đã thêm
    # Đọc nội dung → trang/slide → chunk → embed → add
    if kind == "pdf":
        pages = read_pdf(local_path)
    else:
        pages = read_pptx(local_path)

    all_chunks, map_info = [], []  # map_info: (page, chunk_id, text)
    for page_no, page_text in pages:
        for cid, ch in enumerate(chunk_text(page_text)):
            all_chunks.append(ch)
            map_info.append((page_no, cid, ch))

    vecs = embed_texts(client, all_chunks)
    vecs = l2_normalize(vecs)

    if vecs.shape[0] > 0:
        # Nếu index là FlatIP, thêm trực tiếp
        index.add(vecs)

    # Cập nhật meta: lưu dải vị trí vector theo file để truy nguồn
    total = index.ntotal
    start = total - vecs.shape[0]
    meta["files"][file_id] = {
        "title": title,
        "md5": md5_of_file(local_path),
        "modified": os.path.getmtime(local_path),  # dùng timestamp local tải về
        "ranges": [(start, total)],               # có thể tách nhiều đợt nếu cần
        "granularity": [{"page": p, "chunk": c} for (p, c, _) in map_info]
    }
    return vecs.shape[0]

def md5_of_file(path: str) -> str:
    with open(path, "rb") as f:
        return md5_bytes(f.read())
