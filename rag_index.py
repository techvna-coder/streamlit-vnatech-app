import os, pickle, numpy as np
from openai import OpenAI
import tiktoken
from pypdf import PdfReader
from pptx import Presentation

EMBED_MODEL = "text-embedding-3-small"
MAX_TOKENS, OVERLAP = 400, 80
TOP_K = 10
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

def read_pdf(path):
    reader = PdfReader(path)
    return "\n".join([p.extract_text() or "" for p in reader.pages])

def read_pptx(path):
    prs = Presentation(path)
    texts = []
    for i, slide in enumerate(prs.slides):
        parts = [sh.text for sh in slide.shapes if hasattr(sh,"text")]
        texts.append(f"[Slide {i+1}] " + " ".join(parts))
    return "\n".join(texts)

def chunk_text(text):
    enc = tiktoken.get_encoding("cl100k_base")
    toks = enc.encode(text)
    chunks, start = [], 0
    while start < len(toks):
        end = min(start + MAX_TOKENS, len(toks))
        chunks.append(enc.decode(toks[start:end]))
        if end == len(toks): break
        start = end - OVERLAP
    return chunks

def embed_texts(texts):
    resp = client.embeddings.create(model=EMBED_MODEL, input=texts)
    return np.array([d.embedding for d in resp.data])

def build_or_load(files, cache="index.pkl"):
    if os.path.exists(cache):
        with open(cache,"rb") as f: return pickle.load(f)

    chunks, metas = [], []
    for f in files:
        if f.endswith(".pdf"): content = read_pdf(f)
        elif f.endswith(".pptx"): content = read_pptx(f)
        else: continue
        for i, c in enumerate(chunk_text(content)):
            chunks.append(c); metas.append({"file": f, "id": i})

    embeds = embed_texts(chunks)
    index = {"chunks": chunks, "metas": metas, "embeds": embeds}
    with open(cache,"wb") as f: pickle.dump(index,f)
    return index

def retrieve(query, index, k=TOP_K):
    qvec = embed_texts([query])[0]
    sims = index["embeds"] @ qvec / (
        np.linalg.norm(index["embeds"],axis=1)*np.linalg.norm(qvec))
    top = np.argsort(-sims)[:k]
    return [(float(sims[i]), index["metas"][i], index["chunks"][i]) for i in top]
